from flask import Flask, request, jsonify, render_template, redirect
from together import Together
from flask_cors import CORS
from werkzeug.utils import secure_filename
from flask_login import LoginManager, login_user, login_required, logout_user, current_user
from werkzeug.security import generate_password_hash, check_password_hash
from dotenv import load_dotenv
import markdown
import os
import PyPDF2
import docx
import pandas as pd
from datetime import datetime
from models import db, User, Message, Conversation

# Novos imports para PPTX e imagens
from pptx import Presentation
from PIL import Image
import pytesseract

# Configuração do Tesseract no Windows (ajuste o caminho se necessário)
pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'

# Carregar variáveis de ambiente
load_dotenv()

# Caminho base do projeto
basedir = os.path.abspath(os.path.dirname(__file__))

# Garante que a pasta instance exista
instance_path = os.path.join(basedir, 'instance')
os.makedirs(instance_path, exist_ok=True)

# Inicialização do Flask
app = Flask(__name__, static_folder="static", template_folder="Templates", instance_path=instance_path)

# Configurações
app.config['SQLALCHEMY_DATABASE_URI'] = f"sqlite:///{os.path.join(instance_path, 'db.sqlite3')}"
app.config['SECRET_KEY'] = os.getenv("SECRET_KEY", "chave-padrao")
db.init_app(app)
CORS(app)

# Login
login_manager = LoginManager()
login_manager.login_view = 'login_page'
login_manager.init_app(app)

@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

# Cliente da IA
client = Together(api_key=os.getenv("TOGETHER_API_KEY"))

# Lê o conteúdo do prompt do arquivo externo
with open("prompt.txt", "r", encoding="utf-8") as file:
    PROMPT_TEA = file.read()


@app.route("/")
@login_required
def home():
    return render_template("index.html")

@app.route("/chat", methods=["POST"])
@login_required
def chat():
    try:
        data = request.get_json()
        user_message = data.get("message", "")
        conversation_id = data.get("conversation_id")

        if not user_message:
            return jsonify({"error": "Mensagem vazia"}), 400

        if conversation_id:
            conversation = Conversation.query.get(conversation_id)
            if not conversation or conversation.user_id != current_user.id:
                return jsonify({"error": "Conversa inválida"}), 403
        else:
            conversation = Conversation(user_id=current_user.id, title="Nova conversa")
            db.session.add(conversation)
            db.session.commit()

        response = client.chat.completions.create(
            model="meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
            messages=[{
                "role": "user",
                "content": PROMPT_TEA + user_message
            }]
        )

        raw_message = response.choices[0].message.content
        bot_message = markdown.markdown(raw_message)

        mensagem = Message(
            user_id=current_user.id,
            conversation_id=conversation.id,
            user_message=user_message,
            bot_response=raw_message
        )
        db.session.add(mensagem)
        db.session.commit()

        return jsonify({"response": bot_message, "conversation_id": conversation.id})

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/upload", methods=["POST"])
@login_required
def upload():
    file = request.files.get("file")
    if not file:
        return jsonify({"error": "Arquivo não enviado"}), 400

    filename = secure_filename(file.filename)
    ext = filename.rsplit('.', 1)[-1].lower()

    try:
        if ext == "txt":
            content = file.read().decode("utf-8")
        elif ext == "pdf":
            reader = PyPDF2.PdfReader(file)
            content = "\n".join(page.extract_text() or "" for page in reader.pages)
        elif ext == "docx":
            doc = docx.Document(file)
            content = "\n".join(p.text for p in doc.paragraphs)
        elif ext in ["xls", "xlsx"]:
            df = pd.read_excel(file)
            content = df.to_string(index=False)
        elif ext == "pptx":
            prs = Presentation(file)
            content = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif ext in ["jpg", "jpeg", "png"]:
            image = Image.open(file.stream)
            content = pytesseract.image_to_string(image)
        else:
            return jsonify({"error": "Tipo de arquivo não suportado."}), 400

        response = client.chat.completions.create(
            model="meta-llama/Llama-3.3-70B-Instruct-Turbo-Free",
            messages=[{
                "role": "user",
                "content": PROMPT_TEA + content
            }]
        )

        raw_message = response.choices[0].message.content
        bot_message = markdown.markdown(raw_message)

        mensagem = Message(user_id=current_user.id, user_message=f"[Arquivo: {filename}]", bot_response=raw_message)
        db.session.add(mensagem)
        db.session.commit()

        return jsonify({"response": bot_message})

    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route("/login", methods=["GET"])
def login_page():
    return render_template("login.html")

@app.route("/login", methods=["POST"])
def login():
    data = request.form if request.form else request.get_json()
    username = data.get("username")
    password = data.get("password")

    user = User.query.filter_by(username=username).first()

    if user and check_password_hash(user.password, password):
        login_user(user)
        return redirect("/") if request.form else jsonify({"message": "Login realizado com sucesso"})

    if request.form:
        return render_template("login.html", error="Credenciais inválidas")
    return jsonify({"error": "Credenciais inválidas"}), 401

@app.route("/register", methods=["GET"])
def register_page():
    return render_template("register.html")

@app.route("/register", methods=["POST"])
def register():
    data = request.form if request.form else request.get_json()
    username = data.get("username")
    password = data.get("password")

    if User.query.filter_by(username=username).first():
        return render_template("register.html", error="Usuário já existe")

    hashed_password = generate_password_hash(password)
    user = User(username=username, password=hashed_password)
    db.session.add(user)
    db.session.commit()

    return redirect("/login")

@app.route("/logout", methods=["POST", "GET"])
@login_required
def logout():
    logout_user()
    return redirect("/login")

@app.route("/nova-conversa", methods=["POST"])
@login_required
def nova_conversa():
    conversa = Conversation(user_id=current_user.id, title="Nova conversa")
    db.session.add(conversa)
    db.session.commit()
    return jsonify({"message": "Conversa criada", "conversa_id": conversa.id})

@app.route("/historico", methods=["GET"])
@login_required
def historico():
    conversas = Conversation.query.filter_by(user_id=current_user.id).order_by(Conversation.created_at.desc()).all()
    data = []
    for c in conversas:
        data.append({
            "id": c.id,
            "title": c.title,
            "created_at": c.created_at.strftime("%d/%m/%Y %H:%M")
        })
    return jsonify(data)

@app.route("/mensagens/<int:conversa_id>")
@login_required
def mensagens(conversa_id):
    conversa = Conversation.query.get_or_404(conversa_id)
    if conversa.user_id != current_user.id:
        return jsonify({"error": "Acesso negado"}), 403

    mensagens = Message.query.filter_by(conversation_id=conversa_id).order_by(Message.timestamp.asc()).all()
    return jsonify([
        {
            "usuario": msg.user_message,
            "assistente": msg.bot_response,
            "timestamp": msg.timestamp.strftime("%H:%M")
        } for msg in mensagens
    ])


# Criar o banco de dados se necessário
if __name__ == "__main__":
    with app.app_context():
        db.create_all()
    app.run(debug=True)
